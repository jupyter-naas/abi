from lib.abi.integration.integration import Integration, IntegrationConfiguration
from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from abi import logger
import json
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL


@dataclass
class PowerPointIntegrationConfiguration(IntegrationConfiguration):
    """Configuration for PowerPoint integration.
    
    Attributes:
        template_path (str, optional): Path to PowerPoint template file
    """
    template_path: Optional[str] = None

class PowerPointIntegration(Integration):
    """PowerPoint integration for creating and modifying presentations.
    
    This class provides methods to interact with PowerPoint files using python-pptx,
    allowing creation and modification of slides, shapes, and content.
    
    Attributes:
        __configuration (PowerPointIntegrationConfiguration): Configuration instance
            containing template settings.
    
    Example:
        >>> config = PowerPointIntegrationConfiguration(
        ...     template_path="path/to/template.pptx"
        ... )
        >>> integration = PowerPointIntegration(config)
    """

    __configuration: PowerPointIntegrationConfiguration

    def __init__(self, configuration: PowerPointIntegrationConfiguration):
        super().__init__(configuration)
        self.__configuration = configuration

    def create_presentation(self, template_path: Optional[str] = None) -> Presentation:
        """Create a new presentation.
        
        Args:
            template_path (str, optional): Path to PowerPoint template file
            
        Returns:
            Presentation: New PowerPoint presentation object
        """
        template_path = template_path or self.__configuration.template_path
        if template_path:
            self.__presentation = Presentation(template_path)
        else:
            self.__presentation = Presentation()
        return self.__presentation
    
    def save_presentation(self, presentation: Presentation, output_path: str) -> None:
        """Save the presentation to a file.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            output_path (str): Path to save presentation to
            
        Example:
            >>> integration.save_presentation(ppt, "output.pptx")
        """
        presentation.save(output_path)
    
    def list_slides(self, presentation: Presentation, text: bool = False) -> List[str]:
        """List all slides in the presentation with their shapes.
        
        Args:
            text (bool, optional): Whether to return slide titles or indices. Defaults to False.
        """
        slides = []
        for i, slide in enumerate(presentation.slides):
            shapes = []
            for shape in slide.shapes:
                data = {
                    "slide_number": i,
                    "shape_name": shape.name if hasattr(shape, 'name') else None,
                    "shape_id": shape.shape_id,
                    "shape_type": shape.shape_type,
                    "text": shape.text.encode('ascii', 'ignore').decode('ascii') if hasattr(shape, 'text') else "",
                    "left": shape.left.cm if hasattr(shape, 'left') else None,
                    "top": shape.top.cm if hasattr(shape, 'top') else None,
                    "width": shape.width.cm if hasattr(shape, 'width') else None,
                    "height": shape.height.cm if hasattr(shape, 'height') else None,
                    "rotation": shape.rotation if hasattr(shape, 'rotation') else None,
                    # "has_chart": shape.has_chart if hasattr(shape, 'has_chart') else None,
                    # "has_table": shape.has_table if hasattr(shape, 'has_table') else None,
                    # "has_text_frame": shape.has_text_frame if hasattr(shape, 'has_text_frame') else None,
                    # "is_placeholder": shape.is_placeholder if hasattr(shape, 'is_placeholder') else None,
                    # "click_action": str(shape.click_action) if hasattr(shape, 'click_action') else None,
                }
                if text and data["text"] == "":
                    continue
                shapes.append(data)
            slides.append(
                {
                    "slide_number": i,
                    "shapes": shapes
                }
            )
        return slides
    
    def add_slide(
        self, 
        presentation: Presentation, 
        layout_index: int = 6
    ) -> Tuple[Presentation, int]:
        """Add a new slide to the presentation.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            layout_index (int, optional): Index of slide layout to use. Defaults to 6 (blank layout).
            
        Returns:
            Tuple[Presentation, int]: Updated presentation and index of new slide
            
        Example:
            >>> ppt, slide_idx = integration.add_slide(ppt, layout_index=6)
        """
        slide_layout = presentation.slide_layouts[layout_index]
        presentation.slides.add_slide(slide_layout)
        return presentation, len(presentation.slides) - 1

    def add_shape(
        self,
        presentation: Presentation,
        slide_index: int,
        shape_type: int,
        left: float,
        top: float,
        width: float,
        height: float,
        text: Optional[str] = None,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        line_color: Optional[Tuple[int, int, int]] = None
    ) -> Presentation:
        """Add a shape to a slide.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to add shape to
            shape_type (int): Type of shape to add
            left (float): Left position in inches
            top (float): Top position in inches
            width (float): Width in inches
            height (float): Height in inches
            text (str, optional): Text to add to shape
            fill_color (Tuple[int, int, int], optional): RGB fill color
            line_color (Tuple[int, int, int], optional): RGB line color
            
        Returns:
            Presentation: Updated presentation
            
        Example:
            >>> ppt = integration.add_shape(
            ...     ppt,
            ...     slide_index=0,
            ...     shape_type=1,
            ...     left=1,
            ...     top=1,
            ...     width=2,
            ...     height=1,
            ...     text="Hello World",
            ...     fill_color=(255, 0, 0),
            ...     line_color=(0, 0, 0)
            ... )
        """
        slide = presentation.slides[slide_index]
        shape = slide.shapes.add_shape(
            shape_type,
            Cm(left),
            Cm(top),
            Cm(width),
            Cm(height)
        )
        
        if text:
            shape.text = text
            shape.text_frame.paragraphs[0].font.name = font_name
            shape.text_frame.paragraphs[0].font.size = Pt(font_size)
            shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*font_color)
            
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)
            
        if line_color:
            shape.line.color.rgb = RGBColor(*line_color)
            
        return presentation
    
    def add_text_box(
        self,
        presentation: Presentation,
        slide_index: int,
        left: float,
        top: float, 
        width: float,
        height: float,
        text: str,
        word_wrap: bool = True,
        align: int = PP_ALIGN.LEFT,
        line_spacing: Optional[float] = None,
        font_name: str = "Arial",
        font_size: int = 12,
        bold: bool = False,
        italic: bool = False,
        underline: bool = False,
        font_color: str = "000000"
    ) -> Presentation:
        """Add a text box to a slide.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to add text box to
            left (float): Left position in cm
            top (float): Top position in cm
            width (float): Width in cm
            height (float): Height in cm
            text (str): Text content
            word_wrap (bool, optional): Enable word wrapping. Defaults to True.
            margin_top (float, optional): Top margin in cm. Defaults to 0.1.
            margin_bottom (float, optional): Bottom margin in cm. Defaults to 0.1.
            margin_left (float, optional): Left margin in cm. Defaults to 0.1.
            margin_right (float, optional): Right margin in cm. Defaults to 0.1.
            align (int, optional): Text alignment. Defaults to PP_ALIGN.LEFT.
            line_spacing (float, optional): Line spacing in points. Defaults to None.
            font_name (str, optional): Font name. Defaults to "Arial".
            font_size (int, optional): Font size in points. Defaults to 12.
            bold (bool, optional): Bold text. Defaults to False.
            italic (bool, optional): Italic text. Defaults to False.
            underline (bool, optional): Underline text. Defaults to False.
            font_color (str, optional): Font color as hex string. Defaults to "000000".
            
        Returns:
            Presentation: Updated presentation
            
        Example:
            >>> ppt = integration.add_text_box(
            ...     ppt,
            ...     slide_index=0,
            ...     left=1,
            ...     top=1,
            ...     width=8,
            ...     height=2,
            ...     text="Hello World",
            ...     font_size=24,
            ...     bold=True,
            ...     align=PP_ALIGN.CENTER
            ... )
        """
        slide = presentation.slides[slide_index]
        
        # Add textbox
        txBox = slide.shapes.add_textbox(
            Cm(left),
            Cm(top),
            Cm(width),
            Cm(height)
        )
        
        # Setup textFrame
        textFrame = txBox.text_frame
        textFrame.word_wrap = word_wrap

        # Manage align
        p = textFrame.paragraphs[0]
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = text

        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.bold = bold
        font.italic = italic
        font.underline = underline
        font.color.rgb = RGBColor.from_string(font_color)
        
        return presentation
    
    def update_shape(
        self,
        presentation: Presentation,
        slide_index: int,
        shape_id: int,
        text: Optional[str] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        line_color: Optional[Tuple[int, int, int]] = None,
        left: Optional[float] = None,
        top: Optional[float] = None,
        width: Optional[float] = None,
        height: Optional[float] = None
    ) -> Presentation:
        """Update an existing shape on a slide.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide containing shape
            shape_id (int): ID of shape to update
            text (str, optional): New text for shape
            fill_color (Tuple[int, int, int], optional): New RGB fill color
            line_color (Tuple[int, int, int], optional): New RGB line color
            left (float, optional): New left position in inches
            top (float, optional): New top position in inches
            width (float, optional): New width in inches
            height (float, optional): New height in inches
            
        Returns:
            Presentation: Updated presentation
            
        Example:
            >>> ppt = integration.update_shape(
            ...     ppt,
            ...     slide_index=0,
            ...     shape_id=1,
            ...     text="Updated Text",
            ...     fill_color=(0, 255, 0),
            ...     left=2,
            ...     top=2
            ... )
        """
        slide = presentation.slides[slide_index]
        shape = None
        
        # Find shape by ID
        for s in slide.shapes:
            if s.shape_id == shape_id:
                shape = s
                break
                
        if not shape:
            raise ValueError(f"Shape with ID {shape_id} not found on slide {slide_index}")
            
        if text is not None:
            shape.text = text
            # shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            # shape.text_frame.word_wrap = True
            
        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)
            
        if line_color is not None:
            shape.line.color.rgb = RGBColor(*line_color)
            
        if left is not None:
            shape.left = Inches(left)
            
        if top is not None:
            shape.top = Inches(top)
            
        if width is not None:
            shape.width = Inches(width)
            
        if height is not None:
            shape.height = Inches(height)
            
        return presentation

    def add_image(
        self,
        presentation: Presentation,
        slide_index: int,
        image_path: str,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None
    ) -> Presentation:
        """Add an image to a slide.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to add image to
            image_path (str): Path to image file
            left (float): Left position in inches
            top (float): Top position in inches
            width (float, optional): Width in inches. If None, maintains aspect ratio.
            height (float, optional): Height in inches. If None, maintains aspect ratio.
            
        Returns:
            Presentation: Updated presentation
            
        Example:
            >>> ppt = integration.add_image(
            ...     ppt,
            ...     slide_index=0,
            ...     image_path="path/to/image.png",
            ...     left=1,
            ...     top=1,
            ...     width=2,
            ...     height=2
            ... )
        """
        slide = presentation.slides[slide_index]
        slide.shapes.add_picture(
            image_path,
            Cm(left),
            Cm(top),
            Cm(width) if width else None,
            Cm(height) if height else None
        )
        return presentation

    def add_table(
        self,
        presentation: Presentation,
        slide_index: int,
        rows: int,
        cols: int,
        left: float,
        top: float,
        width: float,
        height: float,
        data: Optional[List[List[str]]] = None
    ) -> Presentation:
        """Add a table to a slide.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to add table to
            rows (int): Number of rows
            cols (int): Number of columns
            left (float): Left position in inches
            top (float): Top position in inches
            width (float): Width in inches
            height (float): Height in inches
            data (List[List[str]], optional): Table data
            
        Returns:
            Presentation: Updated presentation
            
        Example:
            >>> data = [["Header 1", "Header 2"], ["Cell 1", "Cell 2"]]
            >>> ppt = integration.add_table(
            ...     ppt,
            ...     slide_index=0,
            ...     rows=2,
            ...     cols=2,
            ...     left=1,
            ...     top=1,
            ...     width=6,
            ...     height=2,
            ...     data=data
            ... )
        """
        slide = presentation.slides[slide_index]
        table = slide.shapes.add_table(
            rows,
            cols,
            Cm(left),
            Cm(top),
            Cm(width),
            Cm(height)
        ).table
        
        if data:
            for row_idx, row in enumerate(data):
                for col_idx, cell_text in enumerate(row):
                    if row_idx < rows and col_idx < cols:
                        table.cell(row_idx, col_idx).text = str(cell_text)
                        
        return presentation

    def get_presentation_bytes(self, presentation: Presentation) -> bytes:
        """Get the presentation as bytes.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            
        Returns:
            bytes: Presentation file as bytes
            
        Example:
            >>> ppt_bytes = integration.get_presentation_bytes(ppt)
        """
        ppt_stream = io.BytesIO()
        presentation.save(ppt_stream)
        return ppt_stream.getvalue()

    def set_slide_format(
        self,
        presentation: Presentation,
        slide_index: int,
        background_color: Optional[Tuple[int, int, int]] = None,
        title: Optional[str] = None,
        subtitle: Optional[str] = None
    ) -> Presentation:
        """Set formatting for a slide.
        
        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to format
            background_color (Tuple[int, int, int], optional): RGB background color
            title (str, optional): Slide title
            subtitle (str, optional): Slide subtitle
            
        Returns:
            Presentation: Updated presentation
            
        Example:
            >>> ppt = integration.set_slide_format(
            ...     ppt,
            ...     slide_index=0,
            ...     background_color=(255, 255, 255),
            ...     title="My Slide",
            ...     subtitle="Subtitle"
            ... )
        """
        slide = presentation.slides[slide_index]
        
        if background_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*background_color)
            
        if title and hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title'):
            slide.shapes.title.text = title
            
        if subtitle and len(slide.placeholders) > 1:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle
            
        return presentation
    
def as_tools(configuration: PowerPointIntegrationConfiguration):
    """Convert PowerPoint integration into LangChain tools."""
    from langchain_core.tools import StructuredTool
    from pydantic import BaseModel, Field
    
    integration = PowerPointIntegration(configuration)

    class CreatePresentationSchema(BaseModel):
        pass

    class ListSlidesSchema(BaseModel):
        pass

    return [
        StructuredTool(
            name="powerpoint_create_presentation",
            description="Create a new presentation",
            func=lambda **kwargs: integration.create_presentation(**kwargs),
            args_schema=CreatePresentationSchema
        ),
        StructuredTool(
            name="powerpoint_list_slides",
            description="List all slides in the presentation",
            func=lambda **kwargs: integration.list_slides(**kwargs),
            args_schema=ListSlidesSchema
        )
    ]
