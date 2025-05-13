from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
from pathlib import Path
import sys

# Add the project root to the Python path
project_root = Path(__file__).parent.parent.parent.parent.parent
sys.path.append(str(project_root))

def generate_management_deck():
    """Generate a 3-slide PowerPoint deck for management presentation"""
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "GPT-4o Module Overview"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
    
    # Slide 2: Key Features
    features_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = features_slide.shapes.title
    content = features_slide.placeholders[1]
    
    title.text = "Key Features & Benefits"
    
    tf = content.text_frame
    tf.text = "• Advanced multimodal capabilities for enhanced AI interactions\n"
    tf.text += "• Automated conversation analysis and knowledge extraction\n"
    tf.text += "• Integrated document management system\n"
    tf.text += "• Real-time analytics and visualization tools\n"
    tf.text += "• BFO-compliant semantic modeling"
    
    # Slide 3: Usage Statistics
    stats_slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = stats_slide.shapes.title
    content = stats_slide.placeholders[1]
    
    title.text = "Current Usage & Impact"
    
    # Load message data for statistics
    messages_dir = project_root / 'storage' / 'datastore' / 'gpt-4o' / 'messages'
    message_files = list(messages_dir.glob('*.json'))
    
    tf = content.text_frame
    tf.text = f"• Total Conversations: {len(message_files)}\n"
    tf.text += "• Automated knowledge extraction pipeline\n"
    tf.text += "• Daily analytics generation\n"
    tf.text += "• Integrated with multiple interfaces (Terminal, Streamlit, NaasWorkspace)"
    
    # Save the presentation
    output_dir = project_root / 'storage' / 'datastore' / 'analytics' / 'presentations'
    output_dir.mkdir(parents=True, exist_ok=True)
    output_file = output_dir / f'gpt4o_management_deck_{datetime.now().strftime("%Y%m%d")}.pptx'
    prs.save(output_file)
    
    return output_file

if __name__ == "__main__":
    deck_file = generate_management_deck()
    print(f"Generated management deck: {deck_file}")
