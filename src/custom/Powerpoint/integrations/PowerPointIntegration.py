from lib.abi.integration.integration import Integration, IntegrationConfiguration
from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from abi import logger
import json
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL

LOGO_URL = "https://upload.wikimedia.org/wikipedia/commons/thumb/0/0d/Microsoft_Office_PowerPoint_%282019%E2%80%93present%29.svg/2203px-Microsoft_Office_PowerPoint_%282019%E2%80%93present%29.svg.png"

@dataclass
class PowerPointIntegrationConfiguration(IntegrationConfiguration):
    """Configuration for PowerPoint integration.
    
    Attributes:
        template_path (str, optional): Path to PowerPoint template file
    """
    template_path: Optional[str] = None

class PowerPointIntegration(Integration):
    """PowerPoint integration for creating and modifying presentations.
    
    This class provides methods to interact with PowerPoint files using python-pptx,
    allowing creation and modification of slides, shapes, and content.
    
    Attributes:
        __configuration (PowerPointIntegrationConfiguration): Configuration instance
            containing template settings.
    
    Example:
        >>> config = PowerPointIntegrationConfiguration(
        ...     template_path="path/to/template.pptx"
        ... )
        >>> integration = PowerPointIntegration(config)
    """

    __configuration: PowerPointIntegrationConfiguration

    def __init__(self, configuration: PowerPointIntegrationConfiguration):
        super().__init__(configuration)
        self.__configuration = configuration
        
    def create_presentation(self, template_path=None) -> Presentation:
        """Create a new presentation, optionally based on a template.
        
        Args:
            template_path (str, optional): Path to a PowerPoint template file.
                If None, the one specified in configuration will be used.
                If both are None, a blank presentation will be created.
                
        Returns:
            Presentation: A python-pptx Presentation object
        """
        path = template_path or self.__configuration.template_path
        if path:
            return Presentation(path)
        else:
            return Presentation()
    
    def list_slides(self, presentation: Presentation, text: bool = False) -> List[Dict]:
        """List all slides in the presentation with their key properties.
        
        Args:
            presentation (Presentation): The presentation to list slides from
            text (bool, optional): Whether to include text content of shapes. Default is False
            
        Returns:
            List[Dict]: List of dictionaries with slide information
        """
        slides = []
        for i, slide in enumerate(presentation.slides):
            slide_info = {
                "slide_number": i + 1,
                "slide_id": slide.slide_id,
                "shapes": []
            }
            
            # Include shapes if text is True
            if text:
                for j, shape in enumerate(slide.shapes):
                    shape_info = {
                        "shape_id": j + 1,
                        "shape_type": str(shape.shape_type),
                        "text": self.__get_shape_text(shape)
                    }
                    slide_info["shapes"].append(shape_info)
            
            slides.append(slide_info)
        
        return slides
    
    def __get_shape_text(self, shape) -> str:
        """Extract text from a shape if it has any.
        
        Args:
            shape: A PowerPoint shape object
            
        Returns:
            str: Text content of the shape, or empty string if no text
        """
        if hasattr(shape, "text"):
            return shape.text
        return ""
    
    def add_slide(self, presentation: Presentation, layout_index: int = 0) -> Dict:
        """Add a new slide to the presentation.
        
        Args:
            presentation (Presentation): The presentation to add a slide to
            layout_index (int, optional): Index of the slide layout to use. Default is 0
            
        Returns:
            Dict: Information about the new slide
        """
        try:
            layout = presentation.slide_layouts[layout_index]
            slide = presentation.slides.add_slide(layout)
            return {
                "slide_number": len(presentation.slides),
                "slide_id": slide.slide_id,
                "success": True
            }
        except Exception as e:
            logger.error(f"Error adding slide: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def add_text_box(self, slide, text: str, left: float = 1, top: float = 1, 
                    width: float = 8, height: float = 1, 
                    font_size: int = 18, bold: bool = False, 
                    alignment: str = "left") -> Dict:
        """Add a text box to a slide.
        
        Args:
            slide: The slide to add a text box to
            text (str): The text content
            left (float): Left position in inches
            top (float): Top position in inches
            width (float): Width in inches
            height (float): Height in inches
            font_size (int): Font size in points
            bold (bool): Whether text should be bold
            alignment (str): Text alignment ('left', 'center', 'right')
            
        Returns:
            Dict: Information about the new text box
        """
        try:
            # Convert inches to slide coordinate system
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)
            
            # Add text box
            textbox = slide.shapes.add_textbox(left_inches, top_inches, width_inches, height_inches)
            tf = textbox.text_frame
            tf.text = text
            
            # Set paragraph alignment
            paragraph = tf.paragraphs[0]
            if alignment.lower() == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif alignment.lower() == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT
            
            # Set font properties
            font = paragraph.font
            font.size = Pt(font_size)
            font.bold = bold
            
            return {
                "success": True,
                "shape_id": len(slide.shapes) - 1
            }
            
        except Exception as e:
            logger.error(f"Error adding text box: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def add_image(self, slide, image_path: str, left: float = 1, top: float = 1,
                 width: float = 4, height: float = 3) -> Dict:
        """Add an image to a slide.
        
        Args:
            slide: The slide to add an image to
            image_path (str): Path to the image file
            left (float): Left position in inches
            top (float): Top position in inches
            width (float): Width in inches
            height (float): Height in inches
            
        Returns:
            Dict: Information about the new image
        """
        try:
            # Convert inches to slide coordinate system
            left_inches = Inches(left)
            top_inches = Inches(top)
            width_inches = Inches(width)
            height_inches = Inches(height)
            
            # Add image
            slide.shapes.add_picture(image_path, left_inches, top_inches, width_inches, height_inches)
            
            return {
                "success": True,
                "shape_id": len(slide.shapes) - 1
            }
            
        except Exception as e:
            logger.error(f"Error adding image: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def update_shape_text(self, slide, shape_index: int, text: str) -> Dict:
        """Update the text of a shape.
        
        Args:
            slide: The slide containing the shape
            shape_index (int): Index of the shape to update
            text (str): New text content
            
        Returns:
            Dict: Status of the operation
        """
        try:
            shapes = list(slide.shapes)
            if shape_index >= len(shapes):
                return {
                    "success": False,
                    "error": f"Shape index {shape_index} out of range (max: {len(shapes) - 1})"
                }
            
            shape = shapes[shape_index]
            if not hasattr(shape, "text_frame"):
                return {
                    "success": False,
                    "error": "Shape does not have a text frame"
                }
                
            shape.text_frame.text = text
            return {
                "success": True
            }
            
        except Exception as e:
            logger.error(f"Error updating shape text: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }
    
    def save(self, presentation: Presentation, file_path: str) -> Dict:
        """Save the presentation to a file.
        
        Args:
            presentation (Presentation): The presentation to save
            file_path (str): Path to save the presentation to
            
        Returns:
            Dict: Status of the operation
        """
        try:
            presentation.save(file_path)
            return {
                "success": True,
                "file_path": file_path
            }
        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            return {
                "success": False,
                "error": str(e)
            }

def as_tools(configuration: PowerPointIntegrationConfiguration):
    """Convert PowerPoint integration into LangChain tools."""
    from langchain_core.tools import StructuredTool
    from pydantic import BaseModel, Field
    
    integration = PowerPointIntegration(configuration)

    class CreatePresentationSchema(BaseModel):
        pass

    class ListSlidesSchema(BaseModel):
        pass

    return [
        StructuredTool(
            name="powerpoint_create_presentation",
            description="Create a new presentation",
            func=lambda **kwargs: integration.create_presentation(**kwargs),
            args_schema=CreatePresentationSchema
        ),
        StructuredTool(
            name="powerpoint_list_slides",
            description="List all slides in the presentation",
            func=lambda **kwargs: integration.list_slides(**kwargs),
            args_schema=ListSlidesSchema
        )
    ] 