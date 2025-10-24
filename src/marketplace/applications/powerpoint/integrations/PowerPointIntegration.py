from lib.abi.integration.integration import Integration, IntegrationConfiguration
from dataclasses import dataclass
from typing import List, Optional, Tuple, Any, Dict
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
from abi import logger
from pptx.util import Cm
import re
from copy import deepcopy

@dataclass
class PowerPointIntegrationConfiguration(IntegrationConfiguration):
    """Configuration for PowerPoint integration.

    Attributes:
        template_path (str, optional): Path to PowerPoint template file
    """

    template_path: Optional[str] = None


class PowerPointIntegration(Integration):
    """PowerPoint integration for creating and modifying presentations.

    This class provides methods to interact with PowerPoint files using python-pptx,
    allowing creation and modification of slides, shapes, and content.

    Attributes:
        __configuration (PowerPointIntegrationConfiguration): Configuration instance
            containing template settings.

    Example:
        >>> config = PowerPointIntegrationConfiguration(
        ...     template_path="path/to/template.pptx"
        ... )
        >>> integration = PowerPointIntegration(config)
    """

    __configuration: PowerPointIntegrationConfiguration
    __presentation: Optional[PresentationType] = None

    def __init__(self, configuration: PowerPointIntegrationConfiguration):
        super().__init__(configuration)
        self.__configuration = configuration

    def create_presentation(self, template_path: Optional[str] = None) -> PresentationType:
        """Create a new presentation.

        Args:
            template_path (str, optional): Path to PowerPoint template file

        Returns:
            Presentation: New PowerPoint presentation object
        """
        template_path = template_path or self.__configuration.template_path
        if template_path:
            self.__presentation = Presentation(template_path)
        else:
            self.__presentation = Presentation()
        return self.__presentation

    def save_presentation(self, presentation: PresentationType, output_path: str) -> None:
        """Save the presentation to a file.

        Args:
            presentation (Presentation): PowerPoint presentation object
            output_path (str): Path to save presentation to

        Example:
            >>> integration.save_presentation(ppt, "output.pptx")
        """
        presentation.save(output_path)

    def list_slides(self, presentation: PresentationType, text: bool = False) -> List[dict]:
        """List all slides in the presentation with their shapes.

        Args:
            text (bool, optional): Whether to return slide titles or indices. Defaults to False.
        """
        slides = []
        for i, slide in enumerate(presentation.slides):
            shapes = []
            for shape in slide.shapes:
                logger.debug(f"Shape: {shape}")
                data = {
                    "slide_number": i,
                    "shape_name": shape.name if hasattr(shape, "name") else None,
                    "shape_alt_text": shape._element._nvXxPr.cNvPr.attrib.get(
                        "descr", ""
                    ),
                    "shape_is_decorative": (
                        shape.is_decorative if hasattr(shape, "is_decorative") else None
                    ),
                    "shape_id": shape.shape_id,
                    "shape_type": shape.shape_type,
                    "text": shape.text if hasattr(shape, "text") else "",
                    "left": shape.left.cm if hasattr(shape, "left") else None,
                    "top": shape.top.cm if hasattr(shape, "top") else None,
                    "width": shape.width.cm if hasattr(shape, "width") else None,
                    "height": shape.height.cm if hasattr(shape, "height") else None,
                    "rotation": shape.rotation if hasattr(shape, "rotation") else None,
                }
                if text and data["text"] == "":
                    continue
                shapes.append(data)
            slides.append({"slide_number": i, "shapes": shapes})
        return slides
    
    def get_shapes_from_slide(self, slide_number: int) -> list:
        """Get the shapes of the presentation.
        """
        presentation = self.create_presentation()
        shapes: list = []
        for shape in presentation.slides[slide_number].shapes:
            data = {
                "slide_number": slide_number,
                "shape_id": shape.shape_id,
                "shape_type": shape.shape_type,
                "shape_alt_text": shape._element._nvXxPr.cNvPr.attrib.get(
                    "descr", ""
                ),
                "text": shape.text if hasattr(shape, "text") else "",
                "left": shape.left.cm if hasattr(shape, "left") else None,
                "top": shape.top.cm if hasattr(shape, "top") else None,
                "width": shape.width.cm if hasattr(shape, "width") else None,
                "height": shape.height.cm if hasattr(shape, "height") else None,
                "rotation": shape.rotation if hasattr(shape, "rotation") else None,
            }
            shapes.append(data)
        return shapes

    def get_structure(self) -> List[Dict[str, Any]]:
        """Get the structure of the presentation.

        Args:
            presentation (Presentation): PowerPoint presentation object
        """
        presentation = self.create_presentation()
        slides = []
        for i, slide in enumerate(presentation.slides):
            shapes = []
            for shape in slide.shapes:
                alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                # if alt_text == "":
                #     continue
                data = {
                    "slide_number": i,
                    "shape_id": shape.shape_id,
                    "shape_type": shape.shape_type,
                    "shape_alt_text": shape._element._nvXxPr.cNvPr.attrib.get(
                        "descr", ""
                    ),
                    "text": shape.text if hasattr(shape, "text") else "",
                    "left": shape.left.cm if hasattr(shape, "left") else None,
                    "top": shape.top.cm if hasattr(shape, "top") else None,
                    "width": shape.width.cm if hasattr(shape, "width") else None,
                    "height": shape.height.cm if hasattr(shape, "height") else None,
                    "rotation": shape.rotation if hasattr(shape, "rotation") else None,
                }
                shapes.append(data)
            slides.append({"slide_number": i, "shapes": shapes})
        return slides

    def add_slide(
        self, 
        presentation: Optional[PresentationType] = None,
        layout_index: Optional[int] = 6
    ) -> Tuple[PresentationType, int]:
        """Add a new slide to the presentation.

        Args:
            presentation (Presentation): PowerPoint presentation object
            layout_index (int, optional): Index of slide layout to use. Defaults to 6 (blank layout).

        Returns:
            Tuple[Presentation, int]: Updated presentation and index of new slide

        Example:
            >>> ppt, slide_idx = integration.add_slide(ppt, layout_index=6)
        """
        presentation = self.create_presentation() if presentation is None else presentation
        slide_layout = presentation.slide_layouts[layout_index]
        presentation.slides.add_slide(slide_layout)
        return presentation, len(presentation.slides) - 1

    def add_shape(
        self,
        presentation: PresentationType,
        slide_index: int,
        shape_type: int,
        left: float,
        top: float,
        width: float,
        height: float,
        text: Optional[str] = None,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        line_color: Optional[Tuple[int, int, int]] = None,
    ) -> PresentationType:
        """Add a shape to a slide.

        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to add shape to
            shape_type (int): Type of shape to add
            left (float): Left position in inches
            top (float): Top position in inches
            width (float): Width in inches
            height (float): Height in inches
            text (str, optional): Text to add to shape
            fill_color (Tuple[int, int, int], optional): RGB fill color
            line_color (Tuple[int, int, int], optional): RGB line color

        Returns:
            Presentation: Updated presentation

        Example:
            >>> ppt = integration.add_shape(
            ...     ppt,
            ...     slide_index=0,
            ...     shape_type=1,
            ...     left=1,
            ...     top=1,
            ...     width=2,
            ...     height=1,
            ...     text="Hello World",
            ...     fill_color=(255, 0, 0),
            ...     line_color=(0, 0, 0)
            ... )
        """
        slide = presentation.slides[slide_index]
        shape = slide.shapes.add_shape(
            shape_type, Cm(left), Cm(top), Cm(width), Cm(height)
        )

        if text:
            shape.text = text
            if font_name:
                shape.text_frame.paragraphs[0].font.name = font_name
            if font_size is not None:
                shape.text_frame.paragraphs[0].font.size = Pt(font_size)
            if font_color:
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*font_color)

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)

        if line_color:
            shape.line.color.rgb = RGBColor(*line_color)

        return presentation

    def add_text_box(
        self,
        presentation: PresentationType,
        slide_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        word_wrap: bool = True,
        align: int = PP_ALIGN.LEFT,
        line_spacing: Optional[float] = None,
        font_name: str = "Arial",
        font_size: int = 12,
        bold: bool = False,
        italic: bool = False,
        underline: bool = False,
        font_color: str = "000000",
    ) -> PresentationType:
        """Add a text box to a slide.

        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to add text box to
            left (float): Left position in cm
            top (float): Top position in cm
            width (float): Width in cm
            height (float): Height in cm
            text (str): Text content
            word_wrap (bool, optional): Enable word wrapping. Defaults to True.
            margin_top (float, optional): Top margin in cm. Defaults to 0.1.
            margin_bottom (float, optional): Bottom margin in cm. Defaults to 0.1.
            margin_left (float, optional): Left margin in cm. Defaults to 0.1.
            margin_right (float, optional): Right margin in cm. Defaults to 0.1.
            align (int, optional): Text alignment. Defaults to PP_ALIGN.LEFT.
            line_spacing (float, optional): Line spacing in points. Defaults to None.
            font_name (str, optional): Font name. Defaults to "Arial".
            font_size (int, optional): Font size in points. Defaults to 12.
            bold (bool, optional): Bold text. Defaults to False.
            italic (bool, optional): Italic text. Defaults to False.
            underline (bool, optional): Underline text. Defaults to False.
            font_color (str, optional): Font color as hex string. Defaults to "000000".

        Returns:
            Presentation: Updated presentation

        Example:
            >>> ppt = integration.add_text_box(
            ...     ppt,
            ...     slide_index=0,
            ...     left=1,
            ...     top=1,
            ...     width=8,
            ...     height=2,
            ...     text="Hello World",
            ...     font_size=24,
            ...     bold=True,
            ...     align=PP_ALIGN.CENTER
            ... )
        """
        slide = presentation.slides[slide_index]

        # Add textbox
        txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))

        # Setup textFrame
        textFrame = txBox.text_frame
        textFrame.word_wrap = word_wrap

        # Manage align
        p = textFrame.paragraphs[0]
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = text

        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.bold = bold
        font.italic = italic
        font.underline = underline
        font.color.rgb = RGBColor.from_string(font_color)

        return presentation

    def update_shape(
        self,
        presentation: PresentationType,
        slide_index: int,
        shape_id: int,
        text: Optional[str] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        line_color: Optional[Tuple[int, int, int]] = None,
        left: Optional[float] = None,
        top: Optional[float] = None,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> PresentationType:
        """Update an existing shape on a slide while preserving text formatting."""
        slide = presentation.slides[slide_index]
        shape = None

        for s in slide.shapes:
            if s.shape_id == shape_id:
                shape = s
                break

        if not shape:
            raise ValueError(
                f"Shape with ID {shape_id} not found on slide {slide_index}"
            )

        if text is not None:
            # Regular text: clean the text frame more carefully.
            text_frame = shape.text_frame

            # Save formatting from the first run of the first paragraph, if available.
            saved_formatting = None
            if text_frame.paragraphs and text_frame.paragraphs[0].runs:
                src_run = text_frame.paragraphs[0].runs[0]
                saved_formatting = {
                    "name": src_run.font.name,
                    "size": src_run.font.size,
                    "bold": src_run.font.bold,
                    "italic": src_run.font.italic,
                    "underline": src_run.font.underline,
                }

            # If there is at least one paragraph, clear its content.
            if text_frame.paragraphs:
                first_par = text_frame.paragraphs[0]
                # Remove all runs in the first paragraph.
                for r in list(first_par.runs):
                    first_par._p.remove(r._r)
                first_par.text = ""  # Ensure no text remains.
                first_par.bullet = False
                first_par.level = 0

                # Remove all additional paragraphs.
                for p in list(text_frame.paragraphs[1:]):
                    p._element.getparent().remove(p._element)
            else:
                # If no paragraph exists, create one.
                first_par = text_frame.add_paragraph()
                first_par.bullet = False
                first_par.level = 0

            # Split the text into lines (preserving empty lines).
            lines = text.split("\n")

            # Populate the first paragraph with the first line.
            if lines:
                # # Clear existing runs
                # for r in list(first_par.runs):
                #     first_par._p.remove(r._r)
                # Clean line 0
                if lines[0].startswith("-") or lines[0].startswith(" "):
                    lines[0] = lines[0][1:]
                first_par = self.__apply_formatted_text(first_par, lines[0], saved_formatting)
                first_par.bullet = False
                first_par.level = 0
                first_par.space_after = Pt(6)

                # For subsequent lines, create a new paragraph for each.
                for line in lines[1:]:
                    if line.startswith("-") or line.startswith(" "):
                        line = line[1:]
                    new_par = text_frame.add_paragraph()
                    new_par = self.__apply_formatted_text(new_par, line, saved_formatting)
                    new_par.bullet = False
                    new_par.level = 0
                    new_par.space_after = Pt(6)

        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)

        if line_color is not None:
            shape.line.color.rgb = RGBColor(*line_color)

        if left is not None:
            shape.left = Cm(left)

        if top is not None:
            shape.top = Cm(top)

        if width is not None:
            shape.width = Cm(width)

        if height is not None:
            shape.height = Cm(height)

        return presentation

    def __extract_source_text_and_url(self, match: str) -> Tuple[str, str]:
        # Extract text and URL from markdown link format
        url_match = re.search(r'\((.*?)\)', match)
        url = url_match.group(1) if url_match else ""
        if '(' in url:
            url = url.split('(', 1)[1].rstrip(')')
        link_text_match = re.search(r'\[(.*?)\]', match)
        link_text = link_text_match.group(1) if link_text_match else ""
        link_text = link_text.replace("Source:", "").strip()
        if link_text == "Source":
            if "www." in url:
                link_text = url.split("www.")[1].split("/")[0]
            elif "https://" in url:
                link_text = url.split("https://")[1].split("/")[0]
            else:
                link_text = url.split("/")[0]
        return link_text, url

    def __apply_formatted_text(self, paragraph: Any, text: str, saved_formatting: Optional[Dict[str, Any]]) -> Any:
        """Apply formatted text with bold prefix for 'XXX: ' pattern.

        Args:
            paragraph: Paragraph object to add text to
            text (str): Text to format
            saved_formatting (dict): Saved font formatting settings
        """
        pattern_href = r'\[.*?\]\(.*?\)'
        if ":" in text and "Source:" not in text:
            prefix, rest = text.split(":", 1)
            # Add prefix in bold
            run = paragraph.add_run()
            run.text = prefix.strip() + ": "
            if saved_formatting:
                run.font.name = saved_formatting.get("name")
                run.font.size = saved_formatting.get("size")
                run.font.bold = True  # Force bold for prefix
                run.font.italic = saved_formatting.get("italic")
                run.font.underline = saved_formatting.get("underline")
            
            # Add rest of text without bold
            run = paragraph.add_run()

            # Apply hyperlink   
            if re.findall(pattern_href, rest):
                matches = re.findall(pattern_href, rest)
                for match in matches:
                    # Extract text and URL from markdown link format
                    link_text, url = self.__extract_source_text_and_url(deepcopy(match))
                    
                    # Remove the markdown link from text
                    rest = rest.replace(match, "").replace("()", "")
                    
                    # Add text without hyperlink
                    run = paragraph.add_run()
                    run.text = rest.strip() + " "
                    if saved_formatting:
                        run.font.name = saved_formatting.get("name")
                        run.font.size = saved_formatting.get("size")
                        run.font.bold = saved_formatting.get("bold")
                        run.font.italic = saved_formatting.get("italic")
                        run.font.underline = saved_formatting.get("underline")

                    # Add hyperlink at the end of the text
                    run = paragraph.add_run()
                    run.text = "(" + link_text + ")"
                    hlink = run.hyperlink
                    hlink.address = url
                    if saved_formatting:
                        run.font.size = saved_formatting.get("size")

            else:
                run.text = rest
                if saved_formatting:
                    run.font.name = saved_formatting.get("name")
                    run.font.size = saved_formatting.get("size")
                    run.font.bold = saved_formatting.get("bold")
                    run.font.italic = saved_formatting.get("italic")
                    run.font.underline = saved_formatting.get("underline")
        
        elif re.findall(pattern_href, text):
            matches = re.findall(pattern_href, text)
            for match in matches:
                # Extract text and URL from markdown link format
                link_text, url = self.__extract_source_text_and_url(deepcopy(match))
                
                # Remove the markdown link from text
                text = text.replace(match, "").replace("()", "")
                
                # Add text without hyperlink
                run = paragraph.add_run()
                run.text = text.strip() + " "

                # Add hyperlink at the end of the text
                run = paragraph.add_run()
                run.text = "(" + link_text + ")"
                hlink = run.hyperlink
                hlink.address = url
                if saved_formatting:
                    run.font.size = saved_formatting.get("size")
        else:
            # Handle normal lines
            paragraph.text = text
            if saved_formatting and paragraph.runs:
                run = paragraph.runs[0]
                run.font.name = saved_formatting.get("name")
                run.font.size = saved_formatting.get("size")
                run.font.bold = saved_formatting.get("bold")
                run.font.italic = saved_formatting.get("italic")
                run.font.underline = saved_formatting.get("underline")
        
        return paragraph

    def add_image(
        self,
        presentation: PresentationType,
        slide_index: int,
        image_path: str,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> PresentationType:
        """Add an image to a slide.

        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to add image to
            image_path (str): Path to image file
            left (float): Left position in inches
            top (float): Top position in inches
            width (float, optional): Width in inches. If None, maintains aspect ratio.
            height (float, optional): Height in inches. If None, maintains aspect ratio.

        Returns:
            Presentation: Updated presentation

        Example:
            >>> ppt = integration.add_image(
            ...     ppt,
            ...     slide_index=0,
            ...     image_path="path/to/image.png",
            ...     left=1,
            ...     top=1,
            ...     width=2,
            ...     height=2
            ... )
        """
        slide = presentation.slides[slide_index]
        slide.shapes.add_picture(
            image_path,
            Cm(left),
            Cm(top),
            Cm(width) if width else None,
            Cm(height) if height else None,
        )
        return presentation

    def add_table(
        self,
        presentation: PresentationType,
        slide_index: int,
        rows: int,
        cols: int,
        left: float,
        top: float,
        width: float,
        height: float,
        data: Optional[List[List[str]]] = None,
    ) -> PresentationType:
        """Add a table to a slide.

        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to add table to
            rows (int): Number of rows
            cols (int): Number of columns
            left (float): Left position in inches
            top (float): Top position in inches
            width (float): Width in inches
            height (float): Height in inches
            data (List[List[str]], optional): Table data

        Returns:
            Presentation: Updated presentation

        Example:
            >>> data = [["Header 1", "Header 2"], ["Cell 1", "Cell 2"]]
            >>> ppt = integration.add_table(
            ...     ppt,
            ...     slide_index=0,
            ...     rows=2,
            ...     cols=2,
            ...     left=1,
            ...     top=1,
            ...     width=6,
            ...     height=2,
            ...     data=data
            ... )
        """
        slide = presentation.slides[slide_index]
        table = slide.shapes.add_table(
            rows, cols, Cm(left), Cm(top), Cm(width), Cm(height)
        ).table

        if data:
            for row_idx, row in enumerate(data):
                for col_idx, cell_text in enumerate(row):
                    if row_idx < rows and col_idx < cols:
                        table.cell(row_idx, col_idx).text = str(cell_text)

        return presentation
    
    def replace_table(
        self,
        presentation: PresentationType,
        slide_index: int,
        shape_id: int,
        data: List[List[str]]
    ) -> PresentationType:
        """Replace an existing table while preserving position and size."""
        slide = presentation.slides[slide_index]
        
        table_shape = next(
            (s for s in slide.shapes 
            if s.shape_id == shape_id and s.shape_type == MSO_SHAPE_TYPE.TABLE), 
            None
        )
        if not table_shape:
            raise ValueError(f"Table with ID {shape_id} not found on slide {slide_index}")

        left = table_shape.left.cm
        top = table_shape.top.cm
        width = table_shape.width.cm
        height = table_shape.height.cm

        table_shape._element.getparent().remove(table_shape._element)

        return self.add_table(
            presentation,
            slide_index,
            len(data),
            len(data[0]) if data else 0,
            left=left,
            top=top,
            width=width,
            height=height,
            data=data
        )

    def get_presentation_bytes(self, presentation: PresentationType) -> bytes:
        """Get the presentation as bytes.

        Args:
            presentation (Presentation): PowerPoint presentation object

        Returns:
            bytes: Presentation file as bytes

        Example:
            >>> ppt_bytes = integration.get_presentation_bytes(ppt)
        """
        ppt_stream = io.BytesIO()
        presentation.save(ppt_stream)
        return ppt_stream.getvalue()

    def set_slide_format(
        self,
        presentation: PresentationType,
        slide_index: int,
        background_color: Optional[Tuple[int, int, int]] = None,
        title: Optional[str] = None,
        subtitle: Optional[str] = None,
    ) -> PresentationType:
        """Set formatting for a slide.

        Args:
            presentation (Presentation): PowerPoint presentation object
            slide_index (int): Index of slide to format
            background_color (Tuple[int, int, int], optional): RGB background color
            title (str, optional): Slide title
            subtitle (str, optional): Slide subtitle

        Returns:
            Presentation: Updated presentation

        Example:
            >>> ppt = integration.set_slide_format(
            ...     ppt,
            ...     slide_index=0,
            ...     background_color=(255, 255, 255),
            ...     title="My Slide",
            ...     subtitle="Subtitle"
            ... )
        """
        slide = presentation.slides[slide_index]

        if background_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*background_color)

        if title and hasattr(slide, "shapes") and hasattr(slide.shapes, "title"):
            slide.shapes.title.text = title

        if subtitle and len(slide.placeholders) > 1:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle

        return presentation

    def _get_layout_in_dst(self, dst_prs, src_layout_name: str):
        """
        Return a slide layout from *dst_prs* that best matches the *src_layout_name*.
        Using a layout object from a different Presentation triggers duplicate ZIP part names.
        """
        # Try exact layout name across all masters
        for m in dst_prs.slide_masters:
            for lyt in m.slide_layouts:
                if lyt.name == src_layout_name:
                    return lyt
                    
        # Fallback to a common layout name
        common_names = [
            "Title and Content", 
            "Title Slide", 
            "Section Header", 
            "Two Content",
            "Comparison", 
            "Title Only", 
            "Blank", 
            "Content with Caption",
            "Picture with Caption",
        ]
        for name in common_names:
            for m in dst_prs.slide_masters:
                for lyt in m.slide_layouts:
                    if lyt.name == name:
                        return lyt
                        
        # Final fallback
        return dst_prs.slide_layouts[0]

    def duplicate_slide(
        self,
        source_presentation: PresentationType,
        source_slide_number: int,
        presentation: PresentationType
    ) -> Tuple[PresentationType, int]:
        """Duplicate a slide while keeping the same layout and content.

        Args:
            source_presentation (Presentation): Source PowerPoint presentation object
            source_slide_number (int): Index of the slide to duplicate (0-based)
            presentation (Presentation): Target PowerPoint presentation object to add the duplicated slide to

        Returns:
            Tuple[Presentation, int]: Updated presentation and index of the duplicated slide

        Raises:
            ValueError: If source_slide_number is invalid

        Example:
            >>> ppt, new_slide_idx = integration.duplicate_slide(source_ppt, 0, target_ppt)
            >>> ppt, new_slide_idx = integration.duplicate_slide(source_ppt, 2, target_ppt)
        """
        src_slide = source_presentation.slides[source_slide_number]

        # Keep the same layout
        src_layout_name = src_slide.slide_layout.name
        layout = self._get_layout_in_dst(presentation, src_layout_name)
        new = presentation.slides.add_slide(layout)

        # Clear any shapes that might exist (from layout) - more robust approach
        shapes_to_remove = list(new.shapes)
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                new.shapes._spTree.remove(sp)
            except Exception as e:
                logger.warning(f"Could not remove shape: {e}")

        # Copy all shapes and pictures
        for s in src_slide.shapes:
            shape_id = s.shape_id
            shape_type = s.shape_type
            shape_text = s.text if hasattr(s, "text") else ""
            print(f"Shape ID: {shape_id}, Shape Type: {shape_type}, Text: {shape_text}")
            try:                    
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE or (s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape_text == ""):
                    # Copy pictures at same position/size
                    try:
                        blob = s.image.blob  # Get the blob directly
                        stream = io.BytesIO(blob)
                        stream.seek(0)
                        new.shapes.add_picture(stream, s.left, s.top, s.width, s.height)
                    except Exception as e:
                        logger.warning(f"Warning: Could not copy picture: {e}")
                        try:
                            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                            rect = new.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, s.left, s.top, s.width, s.height)
                            rect.text_frame.text = "Image could not be copied"
                        except Exception:
                            pass
                else:
                    # Copy other shapes using deep copy
                    new.shapes._spTree.insert_element_before(deepcopy(s._element), 'p:extLst')

            except Exception as e:
                logger.warning(f"Skipping shape due to error: {e}")

        # Copy notes
        if src_slide.has_notes_slide:
            new.notes_slide.notes_text_frame.text = src_slide.notes_slide.notes_text_frame.text

        # Return the index of the newly created slide (always at the end)
        duplicated_slide_index = len(presentation.slides) - 1
        return presentation, duplicated_slide_index
    

    def remove_all_slides(self, presentation: PresentationType) -> PresentationType:
        """Remove all slides from the presentation.

        Args:
            presentation (Presentation): PowerPoint presentation object

        Returns:
            Presentation: Updated presentation with no slides

        Example:
            >>> ppt = integration.remove_all_slides(ppt)
        """
        for _ in range(len(presentation.slides)):
            rId = presentation.slides._sldIdLst[0].rId
            presentation.part.drop_rel(rId)
            presentation.slides._sldIdLst.remove(presentation.slides._sldIdLst[0])
        return presentation
    

    def update_notes(self, presentation: PresentationType, slide_number: int, sources: List[str]) -> PresentationType:
        """
        Add sources as formatted bullet lists to slide notes.
        Args:
            presentation: Presentation to update
            slide_number: Slide number to update
            sources: List of sources
        Returns:
            Presentation with updated notes
        """
        try:
            if not hasattr(presentation, "slides") or presentation.slides is None:
                return presentation
            
            slide = presentation.slides[slide_number]
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame

            # Get or create sources section
            sources_header = "Sources:"
            sources_found = False

            # Check existing paragraphs
            for paragraph in text_frame.paragraphs:
                if sources_header in paragraph.text:
                    sources_found = True
                    # Clear existing source items
                    while len(paragraph._element.getnext()) is not None:
                        next_para = paragraph._element.getnext()
                        paragraph._element.getparent().remove(next_para)
                    break

            if not sources_found:
                # Add sources header if not found
                header_para = text_frame.add_paragraph()
                header_para.text = sources_header
                header_para.font.bold = True

            seen = set()
            for source in sources:
                clean_url = source.strip()
                if clean_url in seen:
                    continue
                seen.add(clean_url)
                bullet_para = text_frame.add_paragraph()
                bullet_para.text = f"• {clean_url}"  # Simulated bullet (as notes section does not support bullet list formatting)
                bullet_para.level = 0
        except Exception as e:
            logger.error(f"Failed processing slide {slide_number}: {str(e)}")
        return presentation


def as_tools(configuration: PowerPointIntegrationConfiguration):
    """Convert PowerPoint integration into LangChain tools."""
    from langchain_core.tools import StructuredTool
    from pydantic import BaseModel

    integration = PowerPointIntegration(configuration)

    class GetStructureSchema(BaseModel):
        pass

    return [
        StructuredTool(
            name="powerpoint_get_structure",
            description="Get the structure of the presentation",
            func=lambda **kwargs: integration.get_structure(**kwargs),
            args_schema=GetStructureSchema,
        ),
    ]
